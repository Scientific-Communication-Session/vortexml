"""
Vortex ML — Retrieval-Augmented Generation (RAG).

A small, self-contained RAG stack:

  documents → chunk → embed → vector store → retrieve top-k → ground a local
  (or cloud) LLM on the retrieved context → answer with citations.

Design goals:
  * Works out of the box with NO model downloads: retrieval uses TF-IDF
    (scikit-learn, already a dependency) and generation can fall back to the
    cloud chatbot key. Everything else is an optional, *better* upgrade.
  * Pluggable generation backends (MLX / llama.cpp / Transformers / Ollama /
    cloud) chosen at query time, each lazily imported and availability-checked,
    with human-readable "what is this / use when" copy for the UI.
  * Apple-Silicon first: MLX is the recommended local backend on this M4.

Nothing here imports `app`, so it stays unit-testable on its own.
"""

import os
import json
import pickle
import platform
import shutil
import time
import urllib.request
from importlib.util import find_spec

import numpy as np

RAG_DIR = os.path.join(os.path.dirname(__file__), "uploads", "rag")
os.makedirs(RAG_DIR, exist_ok=True)

_IS_APPLE_SILICON = platform.system() == "Darwin" and platform.machine() == "arm64"

# Caches for lazily-loaded local models, keyed by model id.
_mlx_cache = {}
_hf_cache = {}
_llamacpp_cache = {}


# ─────────────────────────────────────────────────────────
# 1. Document extraction + chunking
# ─────────────────────────────────────────────────────────
SUPPORTED_EXTS = (
    ".txt", ".md", ".markdown", ".text", ".csv", ".tsv", ".json",
    ".html", ".htm", ".pdf", ".docx", ".xlsx", ".xls", ".pptx", ".rtf",
)


def extract_text(filename, raw_bytes):
    """Pull plain text out of an uploaded document for RAG ingestion.

    Handles everything Claude can ingest, converting to plain text. Heavy/optional
    parsers are imported lazily inside each branch and guarded with a friendly
    ValueError if the library is missing. Unknown/binary types fall back to a
    best-effort utf-8 'ignore' decode (the original behaviour).
    """
    import io

    ext = "." + filename.lower().rsplit(".", 1)[-1] if "." in filename else ""

    # Plain-text-ish formats: decode directly. (.csv/.tsv/.json are already
    # human-readable, so feeding the raw text to the chunker keeps every value
    # searchable.)
    if ext in (".txt", ".md", ".markdown", ".text", ".csv", ".tsv", ".json"):
        return raw_bytes.decode("utf-8", "ignore")

    if ext in (".html", ".htm"):
        # Strip tags to readable text. Prefer BeautifulSoup; fall back to a
        # crude regex strip if it (and lxml) aren't installed.
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            import re
            stripped = re.sub(r"<[^>]+>", " ", raw_bytes.decode("utf-8", "ignore"))
            return re.sub(r"\s+\n", "\n", stripped)
        soup = BeautifulSoup(raw_bytes, "html.parser")
        for tag in soup(["script", "style"]):
            tag.decompose()
        return soup.get_text(separator="\n")

    if ext == ".pdf":
        if find_spec("pypdf") is None:
            raise ValueError("PDF support needs the 'pypdf' package "
                             "(pip install pypdf). Upload .txt/.md/.csv instead.")
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(raw_bytes))
        return "\n\n".join((page.extract_text() or "") for page in reader.pages)

    if ext == ".docx":
        try:
            import docx  # python-docx
        except ImportError:
            raise ValueError("Word (.docx) support needs the 'python-docx' "
                             "package (pip install python-docx).")
        document = docx.Document(io.BytesIO(raw_bytes))
        parts = [p.text for p in document.paragraphs if p.text.strip()]
        # Include table cell text so tabular content is searchable too.
        for table in document.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    parts.append("\t".join(cells))
        return "\n".join(parts)

    if ext in (".xlsx", ".xls"):
        # Render every sheet as readable rows. Reuse pandas (already a dep) which
        # uses openpyxl for .xlsx; .xls needs xlrd.
        try:
            import pandas as pd
        except ImportError:
            raise ValueError("Excel support needs 'pandas' and 'openpyxl'.")
        try:
            sheets = pd.read_excel(io.BytesIO(raw_bytes), sheet_name=None)
        except ImportError:
            raise ValueError("Reading this Excel file needs the 'openpyxl' "
                             "(.xlsx) or 'xlrd' (.xls) package.")
        out = []
        for name, df in sheets.items():
            out.append(f"# Sheet: {name}")
            out.append(df.to_csv(index=False))
        return "\n".join(out)

    if ext == ".pptx":
        try:
            from pptx import Presentation  # python-pptx
        except ImportError:
            raise ValueError("PowerPoint (.pptx) support needs the 'python-pptx' "
                             "package (pip install python-pptx).")
        prs = Presentation(io.BytesIO(raw_bytes))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"# Slide {i}")
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    parts.append(shape.text_frame.text)
        return "\n".join(parts)

    if ext == ".rtf":
        try:
            from striprtf.striprtf import rtf_to_text
        except ImportError:
            raise ValueError("RTF support needs the 'striprtf' package "
                             "(pip install striprtf).")
        return rtf_to_text(raw_bytes.decode("utf-8", "ignore"))

    # Unknown/binary type: best-effort decode rather than hard-failing, matching
    # the codebase's lenient default.
    return raw_bytes.decode("utf-8", "ignore")


def chunk_text(text, chunk_size=900, overlap=150):
    """Split text into overlapping character windows, preferring paragraph
    boundaries so chunks stay semantically coherent."""
    text = (text or "").strip()
    if not text:
        return []
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    chunks, buf = [], ""
    for para in paragraphs:
        if len(buf) + len(para) + 2 <= chunk_size:
            buf = f"{buf}\n\n{para}" if buf else para
        else:
            if buf:
                chunks.append(buf)
            if len(para) <= chunk_size:
                buf = para
            else:
                # Hard-wrap an oversized paragraph with overlap.
                start = 0
                while start < len(para):
                    chunks.append(para[start:start + chunk_size])
                    start += chunk_size - overlap
                buf = ""
    if buf:
        chunks.append(buf)

    # Add a sliding overlap between adjacent chunks for retrieval recall.
    if overlap > 0 and len(chunks) > 1:
        stitched = [chunks[0]]
        for i in range(1, len(chunks)):
            tail = chunks[i - 1][-overlap:]
            stitched.append((tail + "\n" + chunks[i]).strip())
        chunks = stitched
    return chunks


# ─────────────────────────────────────────────────────────
# 2. Embedders
# ─────────────────────────────────────────────────────────
EMBEDDERS = {
    "tfidf": {
        "label": "TF-IDF (fast, local, no download)",
        "description": "Classic keyword/term-frequency vectors via scikit-learn. "
                       "Runs instantly with zero downloads — great for getting "
                       "started or for keyword-heavy documents.",
        "needs": None,
    },
    "minilm": {
        "label": "MiniLM embeddings (semantic, local)",
        "description": "all-MiniLM-L6-v2 sentence embeddings (~80 MB). Understands "
                       "meaning, not just keywords — better retrieval quality. "
                       "Needs the 'sentence-transformers' package.",
        "needs": "sentence_transformers",
    },
}


def embedder_available(name):
    needs = EMBEDDERS.get(name, {}).get("needs")
    return needs is None or find_spec(needs) is not None


def _embed_minilm(texts):
    from sentence_transformers import SentenceTransformer
    model = _hf_cache.get("__minilm__")
    if model is None:
        model = SentenceTransformer("sentence-transformers/all-MiniLM-L6-v2")
        _hf_cache["__minilm__"] = model
    arr = np.asarray(model.encode(texts, normalize_embeddings=True), dtype=np.float32)
    return arr


def _l2_normalize(mat):
    norms = np.linalg.norm(mat, axis=1, keepdims=True)
    norms[norms == 0] = 1.0
    return mat / norms


# ─────────────────────────────────────────────────────────
# 3. Per-knowledge-base vector store (persisted under uploads/rag/<id>)
# ─────────────────────────────────────────────────────────
class KBStore:
    """Holds a knowledge base's chunks + embeddings on disk.

    Layout (one dir per KB):
      chunks.json     — [{id, doc_id, doc_name, text}]
      emb.npy         — float32 (n_chunks x d), L2-normalized
      meta.json       — {embedder, model, chunk_size, overlap}
      vectorizer.pkl  — fitted TfidfVectorizer (tfidf embedder only)
    """

    def __init__(self, kb_id):
        # The dir is created lazily by rebuild(); reads tolerate it being absent.
        self.dir = os.path.join(RAG_DIR, str(kb_id))

    # -- persistence helpers --
    def _p(self, name):
        return os.path.join(self.dir, name)

    def load_chunks(self):
        if os.path.exists(self._p("chunks.json")):
            with open(self._p("chunks.json"), encoding="utf-8") as f:
                return json.load(f)
        return []

    def _meta(self):
        if os.path.exists(self._p("meta.json")):
            with open(self._p("meta.json"), encoding="utf-8") as f:
                return json.load(f)
        return {}

    def destroy(self):
        if os.path.isdir(self.dir):
            shutil.rmtree(self.dir, ignore_errors=True)

    # -- ingest --
    def rebuild(self, chunks, embedder="tfidf", chunk_size=900, overlap=150):
        """(Re)compute embeddings for the full chunk list and persist them.

        TF-IDF must refit its vocabulary whenever the corpus changes, so ingest
        always rebuilds from all chunks — simple and correct for these sizes.
        """
        os.makedirs(self.dir, exist_ok=True)
        texts = [c["text"] for c in chunks]
        if not texts:
            E = np.zeros((0, 1), dtype=np.float32)
        elif embedder == "minilm":
            E = _embed_minilm(texts)
        else:
            from sklearn.feature_extraction.text import TfidfVectorizer
            vec = TfidfVectorizer(max_features=4096, stop_words="english")
            mat = vec.fit_transform(texts)  # already L2-normalized rows
            E = mat.toarray().astype(np.float32)
            with open(self._p("vectorizer.pkl"), "wb") as f:
                pickle.dump(vec, f)

        with open(self._p("chunks.json"), "w", encoding="utf-8") as f:
            json.dump(chunks, f)
        np.save(self._p("emb.npy"), E)
        with open(self._p("meta.json"), "w", encoding="utf-8") as f:
            json.dump({"embedder": embedder, "chunk_size": chunk_size,
                       "overlap": overlap, "n_chunks": len(chunks)}, f)

    # -- query --
    def _embed_query(self, query):
        meta = self._meta()
        if meta.get("embedder") == "minilm":
            return _embed_minilm([query])[0]
        with open(self._p("vectorizer.pkl"), "rb") as f:
            vec = pickle.load(f)
        return vec.transform([query]).toarray().astype(np.float32)[0]

    def search(self, query, k=4):
        chunks = self.load_chunks()
        if not chunks or not os.path.exists(self._p("emb.npy")):
            return []
        E = np.load(self._p("emb.npy"))
        if E.shape[0] == 0:
            return []
        q = self._embed_query(query)
        qn = np.linalg.norm(q)
        if qn > 0:
            q = q / qn
        scores = E @ q
        order = np.argsort(-scores)[:k]
        out = []
        for rank, idx in enumerate(order):
            score = float(scores[idx])
            if score <= 0:
                continue
            c = chunks[int(idx)]
            out.append({
                "n": rank + 1,
                "doc_name": c.get("doc_name"),
                "score": round(score, 4),
                "text": c["text"],
            })
        return out


# ─────────────────────────────────────────────────────────
# 4. Generation backends (pluggable, lazily imported)
# ─────────────────────────────────────────────────────────
# Each backend entry carries the copy the UI shows to explain "what is this /
# when should I use it", plus a few recommended models. The default model is a
# ~3-4B instruct model per the brief (Gemma / Llama / Qwen).
BACKENDS = {
    "cloud": {
        "label": "Cloud (Claude)",
        "tagline": "Zero setup — works right now",
        "description": "Generates with Anthropic's Claude over the same API key "
                       "the tutor chatbot uses. Nothing to install or download. "
                       "Best for trying RAG immediately or when you don't want to "
                       "run a model locally.",
        "use_when": "You want it to just work, or your machine can't run a local LLM.",
        "local": False,
        "novice_default": True,
        "default_model": "claude-opus-4-8",
        "models": ["claude-opus-4-8", "claude-sonnet-4-6", "claude-haiku-4-5"],
        # Friendly names shown in the model picker (the API returns these so the
        # UI can label opaque model IDs).
        "model_labels": {
            "claude-opus-4-8": "Claude Opus 4.8 (1M context) — most capable",
            "claude-sonnet-4-6": "Claude Sonnet 4.6 — balanced",
            "claude-haiku-4-5": "Claude Haiku 4.5 — fastest / cheapest",
        },
        "setup": "Set ANTHROPIC_API_KEY in the server's .env.",
    },
    "mlx": {
        "label": "MLX (Apple Silicon)",
        "tagline": "Fastest local option on this Mac",
        "description": "Apple's MLX framework runs the model on the M-series GPU "
                       "with unified memory. The most efficient way to run a local "
                       "4B model (e.g. Gemma 3 4B, 4-bit) on this machine — fully "
                       "offline and private.",
        "use_when": "You're on Apple Silicon and want fast, private, local answers.",
        "local": True,
        "default_model": "mlx-community/gemma-3-4b-it-4bit",
        "models": [
            "mlx-community/gemma-3-4b-it-4bit",
            "mlx-community/Llama-3.2-3B-Instruct-4bit",
            "mlx-community/Qwen2.5-3B-Instruct-4bit",
        ],
        "setup": "pip install mlx-lm  (Apple Silicon only). The model downloads on first use.",
    },
    "llama_cpp": {
        "label": "llama.cpp (GGUF)",
        "tagline": "Cross-platform, runs quantized GGUF",
        "description": "Runs GGUF-quantized models via llama-cpp-python on CPU or "
                       "Metal/CUDA. Extremely portable and memory-light. Point it at "
                       "a GGUF file or a 'repo_id:filename' to auto-download.",
        "use_when": "You want a tiny, portable runtime or are not on Apple Silicon.",
        "local": True,
        "default_model": "bartowski/gemma-2-2b-it-GGUF:gemma-2-2b-it-Q4_K_M.gguf",
        "models": [
            "bartowski/gemma-2-2b-it-GGUF:gemma-2-2b-it-Q4_K_M.gguf",
            "bartowski/Llama-3.2-3B-Instruct-GGUF:Llama-3.2-3B-Instruct-Q4_K_M.gguf",
        ],
        "setup": "pip install llama-cpp-python  (build can take a few minutes).",
    },
    "transformers": {
        "label": "Transformers (PyTorch)",
        "tagline": "Full-precision, any HF model",
        "description": "Hugging Face Transformers running on MPS/CUDA/CPU via "
                       "PyTorch. The most flexible (any model on the Hub) but the "
                       "heaviest in memory; full or half precision rather than "
                       "aggressive quantization.",
        "use_when": "You need a specific HF model or full precision, and have the RAM.",
        "local": True,
        "default_model": "Qwen/Qwen2.5-3B-Instruct",
        "models": [
            "Qwen/Qwen2.5-3B-Instruct",
            "meta-llama/Llama-3.2-3B-Instruct",
            "google/gemma-3-4b-it",
        ],
        "setup": "pip install transformers accelerate  (some models need a HF login).",
    },
    "ollama": {
        "label": "Ollama (local server)",
        "tagline": "Easiest local setup",
        "description": "Talks to a running Ollama server on this machine. Ollama "
                       "manages downloads and serving for you — the friendliest way "
                       "to run local models, no Python deps needed.",
        "use_when": "You already use Ollama, or want the simplest local setup.",
        "local": True,
        "default_model": "gemma2:2b",
        "models": ["gemma2:2b", "llama3.2", "qwen2.5:3b"],
        "setup": "Install Ollama (ollama.com), run `ollama pull gemma2:2b`, keep it running.",
    },
}

OLLAMA_URL = os.environ.get("OLLAMA_URL", "http://127.0.0.1:11434")


def backend_available(key):
    if key == "cloud":
        return bool(os.environ.get("ANTHROPIC_API_KEY"))
    if key == "mlx":
        return _IS_APPLE_SILICON and find_spec("mlx_lm") is not None
    if key == "llama_cpp":
        return find_spec("llama_cpp") is not None
    if key == "transformers":
        return find_spec("transformers") is not None and find_spec("torch") is not None
    if key == "ollama":
        try:
            urllib.request.urlopen(f"{OLLAMA_URL}/api/tags", timeout=0.4)
            return True
        except Exception:
            return False
    return False


# Backends that can do extended thinking / reasoning. For cloud it's always on
# when requested (Claude); for local backends it depends on the model (Qwen3,
# DeepSeek-R1, etc.) — we pass the request through and surface any reasoning the
# model produces.
THINKING_BACKENDS = {"cloud", "mlx", "transformers", "ollama"}


def supports_thinking(backend):
    return backend in THINKING_BACKENDS


def list_backends():
    recommended = "mlx" if backend_available("mlx") else "cloud"
    out = []
    for key, meta in BACKENDS.items():
        out.append({
            "key": key,
            "label": meta["label"],
            "tagline": meta["tagline"],
            "description": meta["description"],
            "use_when": meta["use_when"],
            "local": meta["local"],
            "available": backend_available(key),
            "default_model": meta["default_model"],
            "models": meta["models"],
            "model_labels": meta.get("model_labels", {}),
            "setup": meta["setup"],
            "recommended": key == recommended,
            "supports_thinking": supports_thinking(key),
        })
    return {"backends": out, "recommended": recommended,
            "apple_silicon": _IS_APPLE_SILICON}


def default_backend():
    """Best available backend: prefer fast local MLX, else cloud, else any."""
    for key in ("mlx", "ollama", "llama_cpp", "transformers", "cloud"):
        if backend_available(key):
            return key
    return "cloud"


# ─────────────────────────────────────────────────────────
# Model inventory — what's downloaded on THIS machine
# ─────────────────────────────────────────────────────────
def _hf_cache_models():
    """{repo_id: {size_bytes}} for models in the local Hugging Face cache."""
    try:
        from huggingface_hub import scan_cache_dir
        info = scan_cache_dir()
        return {r.repo_id: {"size_bytes": int(r.size_on_disk)}
                for r in info.repos if r.repo_type == "model"}
    except Exception:
        return {}


def _ollama_models():
    """{name: {size_bytes}} for models pulled into a local Ollama server."""
    try:
        with urllib.request.urlopen(f"{OLLAMA_URL}/api/tags", timeout=0.5) as r:
            data = json.loads(r.read())
        return {m["name"]: {"size_bytes": m.get("size")} for m in data.get("models", [])}
    except Exception:
        return {}


def local_model_inventory():
    """Raw inventory of models stored on this machine, by source."""
    return {"hf": _hf_cache_models(), "ollama": _ollama_models()}


def device_model_catalog(inventory=None):
    """Curated per-backend RAG model list ('variations') annotated with whether
    each is already downloaded on the device, plus its on-disk size.

    Pass a node's reported `inventory` to build the catalog for a remote device;
    omit it to build for this (the central) machine.
    """
    if inventory is None:
        inventory = local_model_inventory()
    hf = inventory.get("hf", {})
    oll = inventory.get("ollama", {})

    def lookup(backend, model_id):
        if backend == "ollama":
            tag = model_id.split(":")[0]
            for name, meta in oll.items():
                if name == model_id or name.split(":")[0] == tag:
                    return True, meta.get("size_bytes")
            return False, None
        # mlx / transformers use the HF repo id directly; llama.cpp uses repo:file
        repo = model_id.split(":")[0]
        if repo in hf:
            return True, hf[repo].get("size_bytes")
        return False, None

    catalog = []
    for key, meta in BACKENDS.items():
        if not meta["local"]:
            continue  # cloud has no local model files
        models = []
        for m in meta["models"]:
            downloaded, size = lookup(key, m)
            models.append({"id": m, "downloaded": downloaded, "size_bytes": size})
        catalog.append({
            "backend": key,
            "label": meta["label"],
            "available": backend_available(key),
            "default_model": meta["default_model"],
            "models": models,
        })
    return catalog


def installed_models(inv=None):
    """Flat lists of every model stored (this machine, or a node's reported
    inventory), for the manager UI."""
    inv = inv or local_model_inventory()
    return {
        "hf": [{"id": k, "size_bytes": v.get("size_bytes")} for k, v in inv["hf"].items()],
        "ollama": [{"id": k, "size_bytes": v.get("size_bytes")} for k, v in inv["ollama"].items()],
    }


def delete_model(backend, model):
    """Remove a downloaded model from THIS machine to reclaim disk. `backend`
    only selects the storage to delete from: 'ollama' uses the Ollama server,
    anything else deletes the repo from the Hugging Face cache (by repo id)."""
    if backend == "ollama":
        req = urllib.request.Request(
            f"{OLLAMA_URL}/api/delete",
            data=json.dumps({"name": model}).encode(),
            headers={"Content-Type": "application/json"}, method="DELETE")
        urllib.request.urlopen(req, timeout=30)
        return
    repo = model.split(":")[0]
    from huggingface_hub import scan_cache_dir
    for r in scan_cache_dir().repos:
        if r.repo_id == repo and r.repo_type == "model":
            # drop the in-memory handle so the files aren't held open, then delete
            _mlx_cache.pop(model, None)
            _hf_cache.pop(model, None)
            _llamacpp_cache.pop(model, None)
            shutil.rmtree(r.repo_path, ignore_errors=True)
            return
    raise ValueError(f"'{model}' is not downloaded on this device.")


class DownloadCancelled(Exception):
    """Raised inside download_model when its cancel_event is set, to abort a
    download cleanly (between files, mid-chunk via the tqdm hook, or mid-stream
    for Ollama). Callers should treat this as a user-initiated cancel, not a
    failure."""
    pass


def _hf_total_bytes(repo_id, filename=None):
    """Best-effort total download size, in bytes, for a HF repo (whole snapshot)
    or a single file. Returns None if it can't be determined (e.g. offline, no
    file metadata) so the caller falls back to an indeterminate bar."""
    try:
        from huggingface_hub import HfApi
        info = HfApi().model_info(repo_id, files_metadata=True)
        if filename is not None:
            for s in info.siblings:
                if s.rfilename == filename:
                    return int(s.size) if s.size else None
            return None
        total = 0
        seen = False
        for s in info.siblings:
            if s.size:
                total += int(s.size)
                seen = True
        return total if seen else None
    except Exception:
        return None


def _make_progress_tqdm(emit_progress, total_bytes, cancel_event):
    """Build a tqdm subclass to hand to snapshot_download/hf_hub_download.

    HF drives one tqdm bar per file; we accumulate *all* bytes across every bar
    into one shared counter so the UI sees a single overall percentage. update(n)
    is where bytes actually flow, so that's where we: (a) honour cancellation
    (raise DownloadCancelled — it propagates out of HF cleanly), and (b) emit
    throttled structured progress with a smoothed speed."""
    try:
        from tqdm.auto import tqdm as _base_tqdm
    except Exception:
        from tqdm import tqdm as _base_tqdm

    state = {
        "downloaded": 0,
        "last_emit_t": 0.0,
        "last_emit_bytes": 0,
        "speed": None,
        "start_t": time.monotonic(),
        "phase": "",
    }

    def _emit(force=False):
        now = time.monotonic()
        dt = now - state["last_emit_t"]
        dbytes = state["downloaded"] - state["last_emit_bytes"]
        # Throttle: at most ~3 emits/sec, OR when ≥1% more has arrived.
        pct_step = (total_bytes and dbytes / total_bytes >= 0.01)
        if not force and dt < 0.33 and not pct_step:
            return
        if dt > 0 and dbytes >= 0:
            inst = dbytes / dt
            # Light exponential smoothing so the number doesn't jitter.
            state["speed"] = inst if state["speed"] is None else 0.6 * state["speed"] + 0.4 * inst
        state["last_emit_t"] = now
        state["last_emit_bytes"] = state["downloaded"]
        percent = (state["downloaded"] / total_bytes * 100) if total_bytes else None
        emit_progress({
            "status": "downloading",
            "phase": state["phase"],
            "downloaded_bytes": state["downloaded"],
            "total_bytes": total_bytes,
            "percent": percent,
            "speed_bps": state["speed"],
        })

    class ProgressTqdm(_base_tqdm):
        def __init__(self, *a, **k):
            super().__init__(*a, **k)
            desc = (k.get("desc") or getattr(self, "desc", "") or "").strip()
            if desc:
                # HF's snapshot-level bar already reads "Fetching N files"; its
                # per-file bars carry the filename (sometimes a path). Use the
                # snapshot desc verbatim, else prefix the filename with "Fetching".
                if desc.lower().startswith("fetching"):
                    state["phase"] = desc
                else:
                    fname = desc.split("/")[-1].strip(":").strip()
                    state["phase"] = f"Fetching {fname}" if fname else state["phase"]

        def update(self, n=1):
            if cancel_event is not None and cancel_event.is_set():
                raise DownloadCancelled()
            if n:
                state["downloaded"] += int(n)
            _emit()
            return super().update(n)

        def close(self):
            _emit(force=True)
            return super().close()

    return ProgressTqdm


def download_model(backend, model, progress=None, cancel_event=None):
    """Download a model onto THIS machine. Blocking; meant to run in a background
    task.

    `progress(payload)` is an optional callback that receives STRUCTURED dicts:
        { "status": "starting"|"downloading"|"complete",
          "phase": <str>, "downloaded_bytes": <int>, "total_bytes": <int|None>,
          "percent": <float 0..100|None>, "speed_bps": <float|None> }
    `cancel_event` is an optional threading.Event; when set the download aborts
    by raising DownloadCancelled. Raises on failure."""
    def emit(payload):
        if progress:
            progress(payload)

    def _check_cancel():
        if cancel_event is not None and cancel_event.is_set():
            raise DownloadCancelled()

    _check_cancel()

    if backend in ("mlx", "transformers", "llama_cpp"):
        from huggingface_hub import snapshot_download, hf_hub_download
        single_file = backend == "llama_cpp" and ":" in model
        if single_file:
            repo, filename = model.split(":", 1)
            total = _hf_total_bytes(repo, filename)
            phase = f"Fetching {filename}"
        else:
            repo, filename = model, None
            total = _hf_total_bytes(model)
            phase = f"Fetching {model}"

        emit({"status": "starting", "phase": phase, "downloaded_bytes": 0,
              "total_bytes": total, "percent": 0.0 if total else None,
              "speed_bps": None})

        tqdm_cls = _make_progress_tqdm(emit, total, cancel_event)
        if single_file:
            hf_hub_download(repo_id=repo, filename=filename, tqdm_class=tqdm_cls)
        else:
            snapshot_download(repo_id=model, tqdm_class=tqdm_cls)

        _check_cancel()
        emit({"status": "complete", "phase": "Download complete.",
              "downloaded_bytes": total or 0, "total_bytes": total,
              "percent": 100.0, "speed_bps": None})

    elif backend == "ollama":
        emit({"status": "starting", "phase": f"Pulling {model} via Ollama…",
              "downloaded_bytes": 0, "total_bytes": None, "percent": None,
              "speed_bps": None})
        payload = json.dumps({"model": model, "stream": True}).encode()
        req = urllib.request.Request(f"{OLLAMA_URL}/api/pull", data=payload,
                                     headers={"Content-Type": "application/json"})
        last_t = time.monotonic()
        last_bytes = 0
        speed = None
        last_emit_t = 0.0
        with urllib.request.urlopen(req, timeout=3600) as resp:
            for raw in resp:
                _check_cancel()
                raw = raw.strip()
                if not raw:
                    continue
                try:
                    line = json.loads(raw)
                except Exception:
                    continue
                if line.get("error"):
                    raise RuntimeError(line["error"])
                total = line.get("total")
                completed = line.get("completed") or 0
                phase = line.get("status") or "Pulling…"
                now = time.monotonic()
                dt = now - last_t
                if dt > 0 and completed >= last_bytes:
                    inst = (completed - last_bytes) / dt
                    speed = inst if speed is None else 0.6 * speed + 0.4 * inst
                last_t, last_bytes = now, completed
                # Throttle to ~3/sec.
                if now - last_emit_t >= 0.33 or completed == total:
                    last_emit_t = now
                    percent = (completed / total * 100) if total else None
                    emit({"status": "downloading", "phase": phase,
                          "downloaded_bytes": int(completed),
                          "total_bytes": int(total) if total else None,
                          "percent": percent,
                          "speed_bps": speed if total else None})
        _check_cancel()
        emit({"status": "complete", "phase": "Pull complete.",
              "downloaded_bytes": last_bytes, "total_bytes": last_bytes or None,
              "percent": 100.0, "speed_bps": None})
    else:
        raise ValueError(f"Backend '{backend}' has no downloadable local model.")


# -- per-backend generation (all guarded; only called when available) --
# Each takes a multi-turn `messages` list and `reasoning` flag, returning
# (text, completion_tokens|None). A None token count is estimated by the caller.
def _mlx_prompt(tokenizer, messages, reasoning):
    if reasoning:
        try:
            return tokenizer.apply_chat_template(messages, add_generation_prompt=True, enable_thinking=True)
        except TypeError:
            pass  # this model's template doesn't support thinking
    return tokenizer.apply_chat_template(messages, add_generation_prompt=True)


# Cloud models that accept adaptive thinking + the `effort` knob: Opus 4.x and
# Sonnet 4.6. Haiku 4.5 supports neither — sending `output_config.effort` (or a
# `thinking` block) to it returns a 400 — so we omit both for it. Opus 4.8 keeps
# the same request surface as 4.7 (adaptive thinking only; no temperature/top_p/
# budget_tokens) and offers a 1M context window natively (no beta header needed).
def _cloud_supports_effort(model):
    m = (model or "").lower()
    if "haiku" in m:
        return False
    return ("opus-4" in m) or ("sonnet-4-6" in m) or ("sonnet-4.6" in m)


def _cloud_kwargs(messages, model, max_tokens, reasoning):
    """Build the messages.create/stream kwargs, gating thinking+effort by model."""
    model = model or BACKENDS["cloud"]["default_model"]
    system = "\n\n".join(m["content"] for m in messages if m["role"] == "system")
    convo = [{"role": m["role"], "content": m["content"]}
             for m in messages if m["role"] in ("user", "assistant")]
    kwargs = dict(model=model, max_tokens=max_tokens, messages=convo)
    if _cloud_supports_effort(model):
        kwargs["thinking"] = {"type": "adaptive"} if reasoning else {"type": "disabled"}
        kwargs["output_config"] = {"effort": "high" if reasoning else "low"}
    if system:
        kwargs["system"] = [{"type": "text", "text": system, "cache_control": {"type": "ephemeral"}}]
    return kwargs


def _chat_cloud(messages, model, temperature, max_tokens, reasoning=False):
    import anthropic
    client = anthropic.Anthropic()
    resp = client.messages.create(**_cloud_kwargs(messages, model, max_tokens, reasoning))
    text = next((b.text for b in resp.content if getattr(b, "type", None) == "text"), "")
    return text, getattr(getattr(resp, "usage", None), "output_tokens", None)


def _chat_mlx(messages, model, temperature, max_tokens, reasoning=False):
    from mlx_lm import load, generate
    key = model or BACKENDS["mlx"]["default_model"]
    if key not in _mlx_cache:
        _mlx_cache[key] = load(key)
    model_obj, tokenizer = _mlx_cache[key]
    prompt = _mlx_prompt(tokenizer, messages, reasoning)
    text = generate(model_obj, tokenizer, prompt=prompt, max_tokens=max_tokens, verbose=False)
    return text, (len(tokenizer.encode(text)) if text else 0)


def _chat_transformers(messages, model, temperature, max_tokens, reasoning=False):
    from transformers import AutoModelForCausalLM, AutoTokenizer
    key = model or BACKENDS["transformers"]["default_model"]
    if key not in _hf_cache:
        tok = AutoTokenizer.from_pretrained(key)
        mdl = AutoModelForCausalLM.from_pretrained(key, torch_dtype="auto", device_map="auto")
        _hf_cache[key] = (mdl, tok)
    mdl, tok = _hf_cache[key]
    try:
        inputs = tok.apply_chat_template(messages, add_generation_prompt=True, return_tensors="pt",
                                         **({"enable_thinking": True} if reasoning else {})).to(mdl.device)
    except TypeError:
        inputs = tok.apply_chat_template(messages, add_generation_prompt=True, return_tensors="pt").to(mdl.device)
    out = mdl.generate(inputs, max_new_tokens=max_tokens, do_sample=temperature > 0,
                       temperature=max(temperature, 0.01))
    gen = out[0][inputs.shape[1]:]
    return tok.decode(gen, skip_special_tokens=True), int(gen.shape[0])


def _chat_llama_cpp(messages, model, temperature, max_tokens, reasoning=False):
    from llama_cpp import Llama
    key = model or BACKENDS["llama_cpp"]["default_model"]
    if key not in _llamacpp_cache:
        if ":" in key and not os.path.exists(key):
            repo, filename = key.split(":", 1)
            llm = Llama.from_pretrained(repo_id=repo, filename=filename, n_ctx=4096, verbose=False)
        else:
            llm = Llama(model_path=key, n_ctx=4096, verbose=False)
        _llamacpp_cache[key] = llm
    llm = _llamacpp_cache[key]
    resp = llm.create_chat_completion(messages=messages, max_tokens=max_tokens, temperature=temperature)
    return resp["choices"][0]["message"]["content"], resp.get("usage", {}).get("completion_tokens")


def _chat_ollama(messages, model, temperature, max_tokens, reasoning=False):
    body = {"model": model or BACKENDS["ollama"]["default_model"], "messages": messages,
            "stream": False, "options": {"temperature": temperature, "num_predict": max_tokens}}
    if reasoning:
        body["think"] = True
    req = urllib.request.Request(f"{OLLAMA_URL}/api/chat", data=json.dumps(body).encode(),
                                 headers={"Content-Type": "application/json"})
    with urllib.request.urlopen(req, timeout=300) as r:
        data = json.loads(r.read())
    return data["message"]["content"], data.get("eval_count")


_CHAT_GENERATORS = {
    "cloud": _chat_cloud, "mlx": _chat_mlx, "transformers": _chat_transformers,
    "llama_cpp": _chat_llama_cpp, "ollama": _chat_ollama,
}


# -- streaming generators: yield (kind, delta) where kind is "reasoning" or
# "answer", as the tokens are produced --
def _stream_cloud(messages, model, temperature, max_tokens, reasoning=False):
    import anthropic
    client = anthropic.Anthropic()
    kwargs = _cloud_kwargs(messages, model, max_tokens, reasoning)
    # Only stream separate reasoning when the model actually supports thinking.
    thinking_on = kwargs.get("thinking", {}).get("type") == "adaptive"
    with client.messages.stream(**kwargs) as stream:
        if not thinking_on:
            for delta in stream.text_stream:
                yield ("answer", delta)
            return
        for event in stream:
            if getattr(event, "type", None) == "content_block_delta":
                d = event.delta
                dt = getattr(d, "type", None)
                if dt == "thinking_delta":
                    yield ("reasoning", getattr(d, "thinking", ""))
                elif dt == "text_delta":
                    yield ("answer", getattr(d, "text", ""))


def _stream_mlx(messages, model, temperature, max_tokens, reasoning=False):
    from mlx_lm import load, stream_generate
    key = model or BACKENDS["mlx"]["default_model"]
    if key not in _mlx_cache:
        _mlx_cache[key] = load(key)
    model_obj, tokenizer = _mlx_cache[key]
    prompt = _mlx_prompt(tokenizer, messages, reasoning)
    for resp in stream_generate(model_obj, tokenizer, prompt, max_tokens=max_tokens):
        yield ("answer", getattr(resp, "text", str(resp)))  # <think> sectioned client-side


def _stream_ollama(messages, model, temperature, max_tokens, reasoning=False):
    body = {"model": model or BACKENDS["ollama"]["default_model"], "messages": messages,
            "stream": True, "options": {"temperature": temperature, "num_predict": max_tokens}}
    if reasoning:
        body["think"] = True
    req = urllib.request.Request(f"{OLLAMA_URL}/api/chat", data=json.dumps(body).encode(),
                                 headers={"Content-Type": "application/json"})
    with urllib.request.urlopen(req, timeout=300) as r:
        for line in r:
            line = line.strip()
            if not line:
                continue
            d = json.loads(line)
            msg = d.get("message") or {}
            if msg.get("thinking"):
                yield ("reasoning", msg["thinking"])
            if msg.get("content"):
                yield ("answer", msg["content"])
            if d.get("done"):
                break


def _stream_llama_cpp(messages, model, temperature, max_tokens, reasoning=False):
    from llama_cpp import Llama
    key = model or BACKENDS["llama_cpp"]["default_model"]
    if key not in _llamacpp_cache:
        if ":" in key and not os.path.exists(key):
            repo, filename = key.split(":", 1)
            _llamacpp_cache[key] = Llama.from_pretrained(repo_id=repo, filename=filename, n_ctx=4096, verbose=False)
        else:
            _llamacpp_cache[key] = Llama(model_path=key, n_ctx=4096, verbose=False)
    for part in _llamacpp_cache[key].create_chat_completion(
            messages=messages, max_tokens=max_tokens, temperature=temperature, stream=True):
        delta = (part["choices"][0].get("delta") or {}).get("content")
        if delta:
            yield ("answer", delta)


_STREAM_GENERATORS = {
    "cloud": _stream_cloud, "mlx": _stream_mlx,
    "ollama": _stream_ollama, "llama_cpp": _stream_llama_cpp,
}


# ─────────────────────────────────────────────────────────
# 5. RAG orchestration
# ─────────────────────────────────────────────────────────
RAG_SYSTEM_PROMPT = (
    "You are a retrieval-augmented assistant. Answer the user's question using "
    "ONLY the numbered context passages provided. Cite the passages you use with "
    "bracketed numbers like [1] or [2]. If the answer is not contained in the "
    "context, say you don't know based on the provided documents — do not invent "
    "facts. Be concise and accurate."
)


def build_context(kb_id, query, top_k=4):
    """Retrieve top-k chunks and assemble the grounded prompt. Returns None when
    nothing relevant was found. Retrieval always happens centrally (the vector
    store lives here) — only the prompt is shipped to a remote device."""
    hits = KBStore(kb_id).search(query, k=top_k)
    if not hits:
        return None
    context = "\n\n".join(f"[{h['n']}] (from {h['doc_name']})\n{h['text']}" for h in hits)
    user = f"Context passages:\n{context}\n\nQuestion: {query}"
    citations = [{"n": h["n"], "doc_name": h["doc_name"], "score": h["score"],
                  "preview": h["text"][:240]} for h in hits]
    return {"system": RAG_SYSTEM_PROMPT, "user": user, "hits": hits, "citations": citations}


def chat(backend, model, messages, temperature=0.7, max_tokens=700, reasoning=False):
    """Multi-turn generation. `messages` is a list of {role, content} with roles
    system/user/assistant. Returns {text, tokens, gen_time, tokens_per_second}.
    Used by the chat endpoint and (via the bundled rag.py) by a node agent."""
    if backend not in _CHAT_GENERATORS:
        raise ValueError(f"Unknown backend: {backend}")
    if not backend_available(backend):
        meta = BACKENDS.get(backend, {})
        raise RuntimeError(f"Backend '{backend}' isn't available. {meta.get('setup', '')}")
    t0 = time.time()
    text, ntok = _CHAT_GENERATORS[backend](messages, model, temperature, max_tokens, reasoning)
    elapsed = time.time() - t0
    text = (text or "").strip()
    if not ntok:
        ntok = max(1, round(len(text) / 4))  # ~4 chars/token estimate
    return {"text": text, "tokens": int(ntok), "gen_time": round(elapsed, 3),
            "tokens_per_second": round(ntok / elapsed, 1) if elapsed > 0 else None}


def generate(backend, model, system, user, temperature=0.2, max_tokens=700):
    """Single-shot convenience wrapper over chat() (used by RAG query)."""
    return chat(backend, model,
                [{"role": "system", "content": system}, {"role": "user", "content": user}],
                temperature, max_tokens)


def stream_chat(backend, model, messages, temperature=0.7, max_tokens=700, reasoning=False):
    """Yield (kind, delta) where kind is "reasoning" or "answer", as tokens are
    produced. Backends without a streaming path (transformers) fall back to one
    final chunk."""
    if backend not in _CHAT_GENERATORS:
        raise ValueError(f"Unknown backend: {backend}")
    if not backend_available(backend):
        meta = BACKENDS.get(backend, {})
        raise RuntimeError(f"Backend '{backend}' isn't available. {meta.get('setup', '')}")
    gen = _STREAM_GENERATORS.get(backend)
    if gen is None:  # non-streaming backend → emit the whole reply at once
        text, _ = _CHAT_GENERATORS[backend](messages, model, temperature, max_tokens, reasoning)
        yield ("answer", text or "")
        return
    yield from gen(messages, model, temperature, max_tokens, reasoning)


def rag_answer(kb_id, query, backend="cloud", model=None, top_k=4,
               temperature=0.2, max_tokens=700):
    """Retrieve, ground, and generate locally. Returns answer + citations +
    chunks + generation stats (tokens/sec)."""
    ctx = build_context(kb_id, query, top_k=top_k)
    if ctx is None:
        return {"answer": "I couldn't find anything relevant in this knowledge "
                          "base. Add documents (or rephrase the question) and try again.",
                "citations": [], "chunks": [], "backend": backend, "model": model,
                "stats": None}
    gen = generate(backend, model, ctx["system"], ctx["user"], temperature, max_tokens)
    return {"answer": gen["text"], "citations": ctx["citations"], "chunks": ctx["hits"],
            "backend": backend, "model": model or BACKENDS[backend]["default_model"],
            "stats": {"tokens": gen["tokens"], "gen_time": gen["gen_time"],
                      "tokens_per_second": gen["tokens_per_second"]}}
